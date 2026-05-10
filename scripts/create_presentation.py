"""Generate a professional PowerPoint deck for the PubMed RAG Chatbot project.

Run from the project root:

    python scripts/create_presentation.py

The output file is:

    presentation/PubMed_RAG_Chatbot_Presentation.pptx

The deck uses editable PowerPoint shapes, text, diagrams, and screenshot
placeholders so it can be customized later without design software.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "presentation"
OUTPUT_FILE = OUTPUT_DIR / "PubMed_RAG_Chatbot_Presentation_Updated.pptx"
ASSET_DIR = PROJECT_ROOT / "assets" / "presentation"
AI_MEDICAL_IMAGE = ASSET_DIR / "ai_medical_visual.png"
TERMINAL_IMAGE = ASSET_DIR / "terminal_ingestion_output.png"
MONGODB_IMAGE = ASSET_DIR / "mongodb_collections_visual.png"
GUI_SCREENSHOT = ASSET_DIR / "gradio_ui_screenshot.png"
OUTPUT_SCREENSHOT = ASSET_DIR / "gradio_output_screenshot.png"

REPO_URL = "https://github.com/anaconda0/rag_pubmed_chatbot"
PROJECT_TITLE = "PubMed RAG Chatbot"
UNIVERSITY_NAME = "Arab Academy for Science, Technology & Maritime Transport"
TEAM_MEMBERS = [
    ("Omar Medhat", "221004675"),
    ("Nagham ElNoshokaty", "221006874"),
    ("Yousef Ashraf", "221005207"),
    ("Marwan Ashraf", "221005698"),
    ("Farouk Faisal", "221006961"),
]

# 16:9 widescreen slide size.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Modern AI / healthcare / cybersecurity inspired palette.
NAVY = RGBColor(7, 15, 34)
NAVY_2 = RGBColor(10, 27, 56)
PANEL = RGBColor(15, 35, 66)
PANEL_2 = RGBColor(18, 49, 89)
CYAN = RGBColor(0, 220, 255)
CYAN_DARK = RGBColor(0, 149, 182)
WHITE = RGBColor(245, 250, 255)
MUTED = RGBColor(176, 197, 220)
GREEN = RGBColor(45, 212, 191)
YELLOW = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
GRAY = RGBColor(93, 116, 145)
BLACK = RGBColor(0, 0, 0)

FONT_TITLE = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Consolas"


def set_fill(shape, color: RGBColor) -> None:
    """Apply a solid fill color to a shape."""

    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor, width: float = 1.0, dash: bool = False) -> None:
    """Style a shape line."""

    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 20,
    color: RGBColor = WHITE,
    bold: bool = False,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> object:
    """Add a text box with consistent typography."""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    return box


def add_rich_text(
    slide,
    lines: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: RGBColor = WHITE,
    bullet: bool = True,
) -> object:
    """Add a concise list of bullet-style lines."""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = FONT_BODY
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(7)
        paragraph.level = 0
        if bullet:
            paragraph.text = f"- {line}"
    return box


def add_title(slide, title: str, subtitle: str = "") -> None:
    """Add a standard title block."""

    add_text(slide, title, 0.62, 0.36, 8.5, 0.55, 28, WHITE, True, FONT_TITLE)
    if subtitle:
        add_text(slide, subtitle, 0.66, 0.94, 8.3, 0.33, 11, CYAN, False, FONT_BODY)


def add_slide_transition(slide) -> None:
    """Add a subtle fade transition to the slide."""

    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))
    slide._element.insert(1, transition)


def add_notes(slide, notes: str) -> None:
    """Add speaker notes for a short 5-minute presentation."""

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def add_background(slide, slide_number: int | None = None) -> None:
    """Create the dark AI-healthcare background used across the deck."""

    bg = slide.shapes.add_shape(SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, NAVY)
    bg.line.fill.background()

    # Cybersecurity-inspired side rail.
    rail = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), SLIDE_H
    )
    set_fill(rail, CYAN)
    rail.line.fill.background()

    # Soft panel glow.
    glow = slide.shapes.add_shape(
        SHAPE.OVAL, Inches(10.1), Inches(-0.9), Inches(4.2), Inches(4.2)
    )
    set_fill(glow, PANEL_2)
    glow.line.fill.background()

    # Network nodes.
    node_points = [
        (10.2, 0.7),
        (11.0, 1.15),
        (12.05, 0.82),
        (11.75, 1.78),
        (12.55, 1.45),
    ]
    for start, end in zip(node_points, node_points[1:]):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start[0]),
            Inches(start[1]),
            Inches(end[0]),
            Inches(end[1]),
        )
        set_line(line, CYAN_DARK, 1.2)
    for x, y in node_points:
        node = slide.shapes.add_shape(
            SHAPE.OVAL, Inches(x - 0.04), Inches(y - 0.04), Inches(0.08), Inches(0.08)
        )
        set_fill(node, CYAN)
        node.line.fill.background()

    # Footer.
    footer = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(0.6), Inches(7.07), Inches(12.0), Inches(0.03)
    )
    set_fill(footer, RGBColor(29, 78, 112))
    footer.line.fill.background()
    add_text(
        slide,
        "PubMed RAG Chatbot | MongoDB + Sentence-Transformers + Gradio",
        0.62,
        7.14,
        6.9,
        0.22,
        8,
        MUTED,
    )
    if slide_number is not None:
        add_text(
            slide,
            f"{slide_number:02d}",
            12.18,
            7.08,
            0.55,
            0.25,
            9,
            CYAN,
            True,
            align=PP_ALIGN.RIGHT,
        )


def add_chip(slide, text: str, x: float, y: float, w: float, color: RGBColor = CYAN) -> None:
    """Add a small labeled chip."""

    chip = slide.shapes.add_shape(
        SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34)
    )
    set_fill(chip, PANEL)
    set_line(chip, color, 1.1)
    add_text(slide, text, x + 0.06, y + 0.07, w - 0.12, 0.18, 8, WHITE, True, align=PP_ALIGN.CENTER)


def add_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str = "",
    border: RGBColor = CYAN_DARK,
    fill: RGBColor = PANEL,
) -> object:
    """Add a dark rounded panel."""

    panel = slide.shapes.add_shape(
        SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(panel, fill)
    set_line(panel, border, 1.3)
    if title:
        add_text(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.28, 13, CYAN, True)
    return panel


def add_icon_card(
    slide,
    label: str,
    detail: str,
    x: float,
    y: float,
    w: float,
    h: float,
    icon_text: str,
    accent: RGBColor = CYAN,
) -> None:
    """Add a reusable card with a simple editable icon."""

    add_panel(slide, x, y, w, h, fill=PANEL)
    circle = slide.shapes.add_shape(
        SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.48), Inches(0.48)
    )
    set_fill(circle, accent)
    circle.line.fill.background()
    add_text(
        slide,
        icon_text,
        x + 0.18,
        y + 0.29,
        0.48,
        0.18,
        8,
        NAVY,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, label, x + 0.78, y + 0.16, w - 0.95, 0.28, 14, WHITE, True)
    add_text(slide, detail, x + 0.78, y + 0.50, w - 0.95, h - 0.58, 10, MUTED)


def add_arrow(slide, x: float, y: float, w: float = 0.55, h: float = 0.18) -> None:
    """Add a compact right arrow between diagram blocks."""

    arrow = slide.shapes.add_shape(SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(arrow, CYAN)
    arrow.line.fill.background()


def add_screenshot_placeholder(
    slide,
    label: str,
    x: float,
    y: float,
    w: float,
    h: float,
    hint: str,
) -> None:
    """Add a dashed screenshot placeholder."""

    box = slide.shapes.add_shape(
        SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(box, NAVY_2)
    set_line(box, CYAN, 1.4, dash=True)
    add_text(slide, label, x + 0.16, y + 0.14, w - 0.32, 0.28, 13, CYAN, True)
    add_text(
        slide,
        hint,
        x + 0.18,
        y + h / 2 - 0.22,
        w - 0.36,
        0.5,
        12,
        MUTED,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    # Small frame lines to make it visually read as a screenshot slot.
    top = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(x + 0.16), Inches(y + 0.58), Inches(w - 0.32), Inches(0.03)
    )
    set_fill(top, CYAN_DARK)
    top.line.fill.background()


def pil_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    """Load a clean Windows font for generated visual assets."""

    font_name = "segoeuib.ttf" if bold else "segoeui.ttf"
    try:
        return ImageFont.truetype(font_name, size)
    except OSError:
        return ImageFont.load_default()


def draw_glow_circle(
    draw: ImageDraw.ImageDraw,
    center: tuple[int, int],
    radius: int,
    color: tuple[int, int, int],
) -> None:
    """Draw layered circles to create a soft glow effect."""

    x, y = center
    for step in range(5, 0, -1):
        alpha_radius = int(radius * step / 2.5)
        shade = tuple(max(0, min(255, value // step + 12)) for value in color)
        draw.ellipse(
            (x - alpha_radius, y - alpha_radius, x + alpha_radius, y + alpha_radius),
            outline=shade,
            width=2,
        )


def create_ai_medical_visual(path: Path) -> None:
    """Create a reusable AI-healthcare visual image for the deck."""

    width, height = 1600, 900
    image = Image.new("RGB", (width, height), (7, 15, 34))
    draw = ImageDraw.Draw(image)

    for y in range(height):
        blend = y / height
        r = int(7 + 7 * blend)
        g = int(15 + 30 * blend)
        b = int(34 + 58 * blend)
        draw.line((0, y, width, y), fill=(r, g, b))

    # Subtle circuit grid.
    for x in range(0, width, 90):
        draw.line((x, 0, x, height), fill=(13, 44, 78), width=1)
    for y in range(0, height, 90):
        draw.line((0, y, width, y), fill=(13, 44, 78), width=1)

    # Medical cross with neural network nodes.
    cx, cy = 1110, 360
    draw_glow_circle(draw, (cx, cy), 155, (0, 220, 255))
    draw.rounded_rectangle((cx - 42, cy - 170, cx + 42, cy + 170), radius=22, fill=(0, 160, 190))
    draw.rounded_rectangle((cx - 170, cy - 42, cx + 170, cy + 42), radius=22, fill=(0, 220, 255))

    nodes = [(170, 190), (320, 135), (465, 250), (250, 365), (490, 445), (360, 610)]
    for start, end in zip(nodes, nodes[1:]):
        draw.line((start, end), fill=(0, 180, 215), width=4)
    for x, y in nodes:
        draw.ellipse((x - 14, y - 14, x + 14, y + 14), fill=(45, 212, 191))
        draw.ellipse((x - 5, y - 5, x + 5, y + 5), fill=(245, 250, 255))

    # Heartbeat line.
    heartbeat = [(125, 710), (260, 710), (300, 650), (345, 770), (405, 610), (465, 710), (760, 710)]
    draw.line(heartbeat, fill=(0, 220, 255), width=7, joint="curve")

    draw.text((110, 95), "AI Healthcare RAG", fill=(245, 250, 255), font=pil_font(58, True))
    draw.text((113, 170), "PubMed abstracts + embeddings + retrieval", fill=(176, 197, 220), font=pil_font(30))
    draw.text((118, 790), "Evidence-grounded answers", fill=(45, 212, 191), font=pil_font(28, True))

    image.save(path)


def create_terminal_visual(path: Path) -> None:
    """Create a terminal-style ingestion screenshot visual."""

    width, height = 1200, 700
    image = Image.new("RGB", (width, height), (8, 16, 32))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((35, 35, width - 35, height - 35), radius=22, fill=(10, 27, 56), outline=(0, 220, 255), width=3)
    draw.rectangle((35, 35, width - 35, 92), fill=(15, 35, 66))
    for index, color in enumerate([(248, 113, 113), (251, 191, 36), (45, 212, 191)]):
        x = 72 + index * 34
        draw.ellipse((x, 55, x + 16, 71), fill=color)

    lines = [
        "$ python src/ingest.py --search obesity --limit 50",
        "Rows matching 'obesity': 664",
        "Detected main text/abstract column: abstractText",
        "Loading sentence-transformers model...",
        "Ingesting rows: 100%|==========| 50/50",
        "Rows read: 50",
        "Chunks created before duplicate checks: 96",
        "Duplicate chunks skipped: 0",
        "New chunks inserted into MongoDB: 96",
    ]
    y = 128
    for index, line in enumerate(lines):
        color = (0, 220, 255) if index == 0 else (245, 250, 255)
        if "inserted" in line or "100%" in line:
            color = (45, 212, 191)
        draw.text((72, y), line, fill=color, font=pil_font(30 if index == 0 else 25))
        y += 56

    image.save(path)


def create_mongodb_visual(path: Path) -> None:
    """Create a MongoDB Compass-style collections visual."""

    width, height = 1200, 700
    image = Image.new("RGB", (width, height), (245, 250, 255))
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, width, 86), fill=(10, 27, 56))
    draw.text((45, 25), "MongoDB Collections", fill=(45, 212, 191), font=pil_font(34, True))
    draw.rounded_rectangle((45, 120, 350, 640), radius=18, fill=(232, 247, 243), outline=(45, 212, 191), width=3)
    draw.text((76, 155), "rag_pubmed_chatbot", fill=(10, 27, 56), font=pil_font(24, True))

    collections = ["pubmed_chunks", "conversation_memory", "indexes"]
    for idx, item in enumerate(collections):
        y = 220 + idx * 78
        fill = (15, 118, 110) if idx == 0 else (255, 255, 255)
        text_color = (255, 255, 255) if idx == 0 else (10, 27, 56)
        draw.rounded_rectangle((72, y, 320, y + 48), radius=10, fill=fill, outline=(143, 211, 201), width=2)
        draw.text((94, y + 12), item, fill=text_color, font=pil_font(19, idx == 0))

    draw.rounded_rectangle((395, 120, 1150, 640), radius=18, fill=(255, 255, 255), outline=(216, 222, 230), width=3)
    draw.text((435, 158), "pubmed_chunks document", fill=(10, 27, 56), font=pil_font(28, True))
    fields = [
        '"text": "medical abstract chunk..."',
        '"embedding": [0.021, -0.113, 0.402, ...]',
        '"chunk_hash": "sha256..."',
        '"metadata": { "pmid", "title", "labels" }',
        '"created_at": ISODate(...)',
    ]
    for idx, field in enumerate(fields):
        y = 230 + idx * 64
        draw.rounded_rectangle((435, y, 1088, y + 40), radius=8, fill=(241, 245, 249), outline=(216, 222, 230), width=1)
        draw.text((458, y + 9), field, fill=(15, 35, 66), font=pil_font(20))

    image.save(path)


def create_visual_assets() -> None:
    """Create generated images used by the presentation."""

    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    create_ai_medical_visual(AI_MEDICAL_IMAGE)
    create_terminal_visual(TERMINAL_IMAGE)
    create_mongodb_visual(MONGODB_IMAGE)


def add_image_fit(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str = "",
) -> bool:
    """Add an image inside a bounded area without stretching it."""

    if not path.exists():
        return False

    add_panel(slide, x, y, w, h, title=title)

    with Image.open(path) as image:
        image_w, image_h = image.size

    image_ratio = image_w / image_h
    box_ratio = w / h

    if image_ratio >= box_ratio:
        pic_w = w - 0.28
        pic_h = pic_w / image_ratio
    else:
        pic_h = h - (0.72 if title else 0.28)
        pic_w = pic_h * image_ratio

    pic_x = x + (w - pic_w) / 2
    pic_y = y + (h - pic_h) / 2 + (0.12 if title else 0)

    slide.shapes.add_picture(
        str(path),
        Inches(pic_x),
        Inches(pic_y),
        width=Inches(pic_w),
        height=Inches(pic_h),
    )
    return True


def add_flow_block(slide, title: str, detail: str, x: float, y: float, w: float, h: float, icon: str) -> None:
    """Add a block for flow diagrams."""

    add_panel(slide, x, y, w, h, fill=PANEL)
    icon_bg = slide.shapes.add_shape(
        SHAPE.OVAL, Inches(x + 0.10), Inches(y + 0.10), Inches(0.44), Inches(0.44)
    )
    set_fill(icon_bg, CYAN)
    icon_bg.line.fill.background()
    add_text(slide, icon, x + 0.11, y + 0.22, 0.42, 0.12, 7, NAVY, True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.64, y + 0.13, w - 0.76, 0.26, 12, WHITE, True)
    add_text(slide, detail, x + 0.15, y + 0.56, w - 0.3, h - 0.62, 9, MUTED)


def add_database_icon(slide, x: float, y: float, w: float, h: float, label: str) -> None:
    """Draw a MongoDB-like database cylinder using native shapes."""

    db = slide.shapes.add_shape(SHAPE.CAN, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(db, PANEL_2)
    set_line(db, GREEN, 1.5)
    add_text(slide, label, x + 0.1, y + h / 2 - 0.12, w - 0.2, 0.22, 11, GREEN, True, align=PP_ALIGN.CENTER)


def add_vector_bars(slide, x: float, y: float) -> None:
    """Draw a simple embedding vector illustration."""

    values = [0.30, 0.75, 0.50, 0.88, 0.42, 0.65, 0.22, 0.58]
    for index, value in enumerate(values):
        bar_h = 0.92 * value
        bar = slide.shapes.add_shape(
            SHAPE.RECTANGLE,
            Inches(x + index * 0.18),
            Inches(y + 0.95 - bar_h),
            Inches(0.10),
            Inches(bar_h),
        )
        set_fill(bar, CYAN if index % 2 else GREEN)
        bar.line.fill.background()
    add_text(slide, "[0.21, 0.78, 0.46, ...]", x - 0.08, y + 1.05, 1.8, 0.2, 8, MUTED, font=FONT_MONO)


def create_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_image_fit(slide, AI_MEDICAL_IMAGE, 8.15, 1.15, 4.25, 3.05)

    add_text(slide, PROJECT_TITLE, 0.72, 1.28, 7.15, 0.75, 39, WHITE, True, FONT_TITLE)
    add_text(
        slide,
        "Retrieval-Augmented Generation for PubMed Medical Abstracts",
        0.78,
        2.10,
        7.2,
        0.35,
        15,
        CYAN,
    )
    add_rich_text(
        slide,
        [
            "MongoDB vector storage",
            "Sentence-transformers embeddings",
            "Conversation memory",
            "Gradio chatbot interface",
        ],
        0.82,
        2.85,
        5.6,
        1.25,
        14,
        WHITE,
    )
    add_panel(slide, 0.82, 4.65, 7.05, 1.92, "Team Members")
    for index, (name, student_id) in enumerate(TEAM_MEMBERS):
        y = 5.05 + index * 0.25
        add_text(slide, name, 1.10, y, 2.90, 0.16, 8, WHITE, True)
        add_text(slide, student_id, 4.18, y, 1.10, 0.16, 8, CYAN, font=FONT_MONO)
    add_text(slide, UNIVERSITY_NAME, 1.10, 6.34, 6.35, 0.16, 8, MUTED)

    add_chip(slide, "AI", 8.45, 4.78, 0.72)
    add_chip(slide, "Healthcare", 9.28, 4.78, 1.15, GREEN)
    add_chip(slide, "Cybersecurity Style", 10.58, 4.78, 1.62)

    add_notes(
        slide,
        "Open by introducing the team and the project as a medical RAG chatbot that uses PubMed abstracts, embeddings, MongoDB, memory, and a Gradio UI.",
    )
    add_slide_transition(slide)


def create_intro(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 2)
    add_title(slide, "Introduction", "Why medical chatbots need grounded answers")

    add_icon_card(
        slide,
        "Chatbot",
        "A software assistant that receives a user question and returns a conversational answer.",
        0.75,
        1.65,
        3.65,
        1.40,
        "CHAT",
    )
    add_icon_card(
        slide,
        "Hallucination Risk",
        "Normal LLMs may generate confident answers that are not supported by the project data.",
        4.85,
        1.65,
        3.65,
        1.40,
        "RISK",
        RED,
    )
    add_icon_card(
        slide,
        "Medical Accuracy",
        "Medical questions require traceable evidence, source chunks, and conservative answers.",
        8.95,
        1.65,
        3.65,
        1.40,
        "MED",
        GREEN,
    )

    add_panel(slide, 1.15, 4.10, 10.85, 1.60, "Project Motivation")
    add_rich_text(
        slide,
        [
            "Answer from medical abstracts, not from unsupported model memory",
            "Show the retrieved PubMed sources used for each answer",
            "Say when the dataset does not contain enough information",
        ],
        1.45,
        4.60,
        10.2,
        0.85,
        16,
        WHITE,
    )

    add_notes(slide, "Explain that the project focuses on reducing unsupported answers in a medical dataset setting.")
    add_slide_transition(slide)


def create_rag_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 3)
    add_title(slide, "What Is RAG?", "Retrieval-Augmented Generation in one flow")

    steps = [
        ("User Question", "Medical query from the chat UI", "Q"),
        ("Retrieval", "Find similar PubMed chunks in MongoDB", "RET"),
        ("Context", "Attach top source chunks to the prompt", "CTX"),
        ("Answer", "Respond only from retrieved context", "ANS"),
    ]
    x = 0.85
    for idx, (title, detail, icon) in enumerate(steps):
        add_flow_block(slide, title, detail, x + idx * 3.05, 2.25, 2.35, 1.55, icon)
        if idx < len(steps) - 1:
            add_arrow(slide, x + idx * 3.05 + 2.45, 2.90, 0.48)

    add_panel(slide, 1.45, 5.05, 10.45, 0.92, "Key Idea")
    add_text(
        slide,
        "The model does not answer from memory alone. It first retrieves project-specific evidence, then builds the final response.",
        1.75,
        5.43,
        9.8,
        0.28,
        15,
        WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, "Use this slide to define RAG simply: retrieval first, answer second.")
    add_slide_transition(slide)


def create_project_idea(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 4)
    add_title(slide, "Project Idea", "A chatbot grounded in PubMed abstracts")

    add_panel(slide, 0.85, 1.55, 5.55, 4.70, "Without RAG")
    add_rich_text(
        slide,
        [
            "Question goes directly to model memory",
            "No project dataset check",
            "Higher risk of unsupported medical claims",
            "Sources are not visible to the user",
        ],
        1.22,
        2.18,
        4.8,
        2.15,
        15,
        MUTED,
    )
    warning = slide.shapes.add_shape(SHAPE.ISOSCELES_TRIANGLE, Inches(2.85), Inches(4.65), Inches(1.0), Inches(0.72))
    set_fill(warning, RED)
    warning.line.fill.background()
    add_text(slide, "!", 3.16, 4.77, 0.35, 0.25, 22, NAVY, True, align=PP_ALIGN.CENTER)

    add_panel(slide, 6.95, 1.55, 5.55, 4.70, "With This Project")
    add_rich_text(
        slide,
        [
            "Question is embedded as a vector",
            "MongoDB returns relevant PubMed chunks",
            "Prompt includes evidence and memory",
            "Answer is generated only from retrieved context",
        ],
        7.32,
        2.18,
        4.8,
        2.15,
        15,
        WHITE,
    )
    check = slide.shapes.add_shape(SHAPE.OVAL, Inches(8.95), Inches(4.62), Inches(0.92), Inches(0.92))
    set_fill(check, GREEN)
    check.line.fill.background()
    add_text(slide, "OK", 9.10, 4.93, 0.60, 0.18, 13, NAVY, True, align=PP_ALIGN.CENTER)

    add_notes(slide, "Contrast a generic chatbot with this grounded RAG workflow.")
    add_slide_transition(slide)


def create_dataset(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 5)
    add_title(slide, "Dataset", "Kaggle PubMed Multilabel Text Classification Dataset")

    add_panel(slide, 0.85, 1.52, 5.55, 4.95, "Dataset Content")
    add_rich_text(
        slide,
        [
            "PubMed medical abstract text",
            "Article titles and PMID identifiers",
            "Multilabel medical categories",
            "MeSH-style disease, treatment, and diagnosis topics",
        ],
        1.20,
        2.05,
        4.95,
        2.35,
        15,
        WHITE,
    )
    add_chip(slide, "Diseases", 1.25, 5.05, 1.0)
    add_chip(slide, "Diagnosis", 2.45, 5.05, 1.05, GREEN)
    add_chip(slide, "Treatments", 3.70, 5.05, 1.18)
    add_chip(slide, "Biology", 5.08, 5.05, 0.95, YELLOW)

    add_panel(slide, 7.05, 1.52, 5.05, 4.95, "Detected CSV Structure")
    add_text(slide, "Title", 7.45, 2.25, 1.1, 0.25, 12, CYAN, True, font=FONT_MONO)
    add_text(slide, "abstractText", 8.70, 2.25, 1.7, 0.25, 12, CYAN, True, font=FONT_MONO)
    add_text(slide, "meshMajor", 10.60, 2.25, 1.25, 0.25, 12, CYAN, True, font=FONT_MONO)
    for row in range(4):
        y = 2.75 + row * 0.58
        row_bg = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(7.35), Inches(y), Inches(4.35), Inches(0.34))
        set_fill(row_bg, NAVY_2)
        set_line(row_bg, RGBColor(41, 83, 118), 0.8)
        add_text(slide, f"PMID {row + 1}", 7.50, y + 0.09, 0.82, 0.13, 7, MUTED, font=FONT_MONO)
        add_text(slide, "medical abstract chunk...", 8.55, y + 0.09, 1.75, 0.13, 7, WHITE, font=FONT_MONO)
        add_text(slide, "labels", 10.72, y + 0.09, 0.65, 0.13, 7, GREEN, font=FONT_MONO)

    add_notes(slide, "Mention that ingestion detects the abstract text column and extracts metadata like title, PMID, and labels.")
    add_slide_transition(slide)


def create_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 6)
    add_title(slide, "System Architecture", "From CSV data to grounded chatbot answers")

    top = [
        ("Dataset CSV", "PubMed abstracts", "CSV"),
        ("Clean + Chunk", "Text preprocessing", "CHK"),
        ("Embeddings", "Sentence-transformers", "VEC"),
        ("MongoDB", "Chunks + vectors", "DB"),
    ]
    bottom = [
        ("Question", "User asks in UI", "Q"),
        ("Retrieval", "Cosine top 5", "TOP"),
        ("Memory", "Short + long term", "MEM"),
        ("Final Answer", "Grounded response", "ANS"),
    ]

    for idx, item in enumerate(top):
        add_flow_block(slide, item[0], item[1], 0.65 + idx * 3.13, 1.55, 2.48, 1.20, item[2])
        if idx < 3:
            add_arrow(slide, 2.95 + idx * 3.13, 2.08, 0.42, 0.16)

    add_database_icon(slide, 9.98, 3.05, 1.5, 1.0, "MongoDB")

    for idx, item in enumerate(bottom):
        add_flow_block(slide, item[0], item[1], 0.65 + idx * 3.13, 5.08, 2.48, 1.20, item[2])
        if idx < 3:
            add_arrow(slide, 2.95 + idx * 3.13, 5.61, 0.42, 0.16)

    # Vertical retrieval connection.
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(10.72), Inches(4.05), Inches(5.05), Inches(5.08)
    )
    set_line(line, CYAN_DARK, 1.2)
    add_text(slide, "Stored vectors searched at question time", 5.75, 3.55, 3.35, 0.28, 10, CYAN, align=PP_ALIGN.CENTER)

    add_notes(slide, "Walk left to right: ingestion pipeline at the top, live chatbot pipeline at the bottom.")
    add_slide_transition(slide)


def create_chunk_embeddings(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 7)
    add_title(slide, "Chunking and Embeddings", "Turning abstracts into searchable vectors")

    add_panel(slide, 0.8, 1.50, 3.2, 4.9, "1. Clean Text")
    add_text(slide, "Remove extra spaces and prepare abstract text for chunking.", 1.10, 2.12, 2.55, 0.80, 14, WHITE)
    add_text(slide, "Abstract text\nTitle metadata\nLabels", 1.15, 3.42, 2.35, 0.85, 13, MUTED, font=FONT_MONO)

    add_arrow(slide, 4.22, 3.55, 0.55, 0.20)

    add_panel(slide, 4.95, 1.50, 3.25, 4.9, "2. Split Into Chunks")
    for index in range(3):
        y = 2.15 + index * 0.72
        chunk = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(y), Inches(2.35), Inches(0.42))
        set_fill(chunk, NAVY_2)
        set_line(chunk, CYAN_DARK, 1)
        add_text(slide, f"Chunk {index}", 5.50, y + 0.12, 1.2, 0.14, 8, WHITE, True, font=FONT_MONO)
    add_text(slide, "Smaller chunks improve retrieval precision.", 5.25, 4.75, 2.55, 0.65, 13, MUTED)

    add_arrow(slide, 8.43, 3.55, 0.55, 0.20)

    add_panel(slide, 9.15, 1.50, 3.25, 4.9, "3. Vector Embeddings")
    add_text(slide, "sentence-transformers converts text meaning into numbers.", 9.45, 2.05, 2.55, 0.72, 13, WHITE)
    add_vector_bars(slide, 9.82, 3.20)
    add_text(slide, "Similar meaning = vectors point in similar directions", 9.42, 5.28, 2.6, 0.50, 11, MUTED)

    add_notes(slide, "Explain chunks as smaller pieces of abstract text, and embeddings as numeric representations of meaning.")
    add_slide_transition(slide)


def create_mongodb_storage(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 8)
    add_title(slide, "MongoDB Storage", "What is stored after ingestion")

    add_database_icon(slide, 0.95, 2.05, 2.25, 3.0, "pubmed_chunks")

    add_panel(slide, 3.70, 1.45, 8.75, 5.15, "Stored Document Shape")
    schema_lines = [
        'text: "medical chunk text..."',
        "embedding: [0.021, -0.113, 0.402, ...]",
        "chunk_hash: duplicate prevention key",
        "metadata.source_file: CSV file name",
        "metadata.row_index / chunk_index",
        "metadata.pmid / title / labels",
    ]
    y = 2.03
    for idx, line in enumerate(schema_lines):
        row = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(4.05), Inches(y + idx * 0.54), Inches(7.75), Inches(0.35))
        set_fill(row, NAVY_2 if idx % 2 == 0 else PANEL_2)
        set_line(row, RGBColor(44, 92, 130), 0.8)
        add_text(slide, line, 4.22, y + idx * 0.54 + 0.10, 7.35, 0.12, 9, WHITE, font=FONT_MONO)

    add_chip(slide, "Unique chunk_hash", 4.10, 5.55, 1.55, GREEN)
    add_chip(slide, "Vector search ready", 5.88, 5.55, 1.55)
    add_chip(slide, "Metadata for sources", 7.66, 5.55, 1.65, YELLOW)
    add_chip(slide, "Labels if available", 9.55, 5.55, 1.48, GREEN)

    add_notes(slide, "Explain that MongoDB stores the text chunks, embeddings, and metadata needed for source display.")
    add_slide_transition(slide)


def create_retrieval(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 9)
    add_title(slide, "Retrieval Process", "Cosine similarity finds the most relevant chunks")

    add_flow_block(slide, "Question Embedding", "Embed the user question with the same model", 0.75, 1.72, 2.7, 1.28, "Q")
    add_arrow(slide, 3.70, 2.25, 0.50, 0.18)
    add_flow_block(slide, "Compare Vectors", "Calculate cosine similarity against MongoDB chunks", 4.45, 1.72, 2.95, 1.28, "COS")
    add_arrow(slide, 7.65, 2.25, 0.50, 0.18)
    add_flow_block(slide, "Top 5 Chunks", "Keep the strongest source chunks for the prompt", 8.40, 1.72, 3.05, 1.28, "5")

    add_panel(slide, 0.95, 4.08, 4.85, 1.75, "Cosine Similarity")
    add_text(slide, "score = dot(question_vector, chunk_vector)", 1.25, 4.75, 4.20, 0.26, 14, WHITE, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "Higher score means closer semantic meaning.", 1.40, 5.22, 3.90, 0.24, 11, MUTED, align=PP_ALIGN.CENTER)

    add_panel(slide, 6.35, 3.70, 5.75, 2.48, "Ranked Results")
    for index, score in enumerate(["0.5358", "0.5312", "0.5134", "0.5015", "0.5004"]):
        y = 4.18 + index * 0.34
        add_text(slide, f"Source {index + 1}", 6.68, y, 0.8, 0.12, 8, WHITE, True, font=FONT_MONO)
        add_text(slide, f"score {score}", 7.72, y, 0.95, 0.12, 8, CYAN, font=FONT_MONO)
        add_text(slide, "metadata + abstract chunk", 8.90, y, 2.65, 0.12, 8, MUTED, font=FONT_MONO)

    add_notes(slide, "Describe retrieval as comparing question meaning against stored chunk meanings and selecting the top five.")
    add_slide_transition(slide)


def create_memory(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 10)
    add_title(slide, "Memory System", "Short-term chat state plus long-term vector memory")

    add_panel(slide, 0.85, 1.45, 5.45, 5.00, "Short-Term Memory")
    add_rich_text(
        slide,
        [
            "Lives inside the active chat session",
            "Stores role/content pairs",
            "Keeps the last 20 messages",
            "Used to understand recent follow-up questions",
        ],
        1.20,
        2.08,
        4.80,
        2.10,
        15,
        WHITE,
    )
    for i, role in enumerate(["user", "assistant", "user"]):
        row = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(1.30), Inches(4.55 + i * 0.36), Inches(4.35), Inches(0.25))
        set_fill(row, NAVY_2)
        set_line(row, CYAN_DARK, 0.8)
        add_text(slide, f"{role}: message content", 1.45, 4.62 + i * 0.36, 3.85, 0.10, 7, MUTED, font=FONT_MONO)

    add_panel(slide, 7.00, 1.45, 5.45, 5.00, "Long-Term Memory")
    add_rich_text(
        slide,
        [
            "Stored in MongoDB",
            "Past user and assistant messages get embeddings",
            "Relevant old memories are retrieved by cosine similarity",
            "Added to prompt as supporting conversation context",
        ],
        7.35,
        2.08,
        4.80,
        2.10,
        15,
        WHITE,
    )
    add_database_icon(slide, 8.82, 4.62, 1.7, 0.95, "memory")

    add_notes(slide, "Explain the difference between session memory and persistent vector memory.")
    add_slide_transition(slide)


def create_gradio_ui(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 11)
    add_title(slide, "Gradio UI", "User-facing chatbot interface")

    if not add_image_fit(slide, GUI_SCREENSHOT, 0.85, 1.45, 7.15, 4.95, "Gradio UI Screenshot"):
        add_screenshot_placeholder(
            slide,
            "Screenshot Placeholder: Gradio UI",
            0.85,
            1.45,
            7.15,
            4.95,
            "Insert screenshot of the running chatbot at http://127.0.0.1:7860",
        )

    add_panel(slide, 8.45, 1.45, 3.95, 4.95, "UI Features")
    add_rich_text(
        slide,
        [
            "Chat interface",
            "Streaming response",
            "Session management",
            "Prompt history panel",
            "Clear and new-session controls",
            "Retrieved source/chunk display",
        ],
        8.78,
        2.05,
        3.30,
        2.95,
        14,
        WHITE,
    )
    add_chip(slide, "Logic stays in src/", 8.78, 5.62, 1.55, GREEN)
    add_chip(slide, "UI stays in ui/", 10.55, 5.62, 1.30)

    add_notes(slide, "Point out that the UI calls the RAG pipeline and does not contain the retrieval logic.")
    add_slide_transition(slide)


def create_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 12)
    add_title(slide, "Demo Flow", "Question -> Retrieval -> Answer")

    if not add_image_fit(slide, OUTPUT_SCREENSHOT, 0.75, 1.35, 7.35, 4.10, "Real Chatbot Output Screenshot"):
        add_screenshot_placeholder(
            slide,
            "Screenshot Placeholder: Chatbot Output",
            0.75,
            1.35,
            7.35,
            4.10,
            "Insert screenshot after asking a question",
        )

    add_image_fit(slide, TERMINAL_IMAGE, 8.45, 1.35, 3.85, 1.80, "Ingestion Output")
    add_image_fit(slide, MONGODB_IMAGE, 8.45, 3.38, 3.85, 2.07, "MongoDB Storage")

    add_panel(slide, 0.95, 5.72, 11.55, 0.72, "Demo Script")
    add_text(slide, 'Question: "What does the dataset say about obesity?"', 1.22, 5.95, 4.80, 0.18, 11, WHITE, True)
    add_text(slide, "Embedding -> cosine retrieval -> top source chunks -> grounded streamed answer", 6.05, 5.95, 5.95, 0.18, 11, CYAN)

    add_notes(slide, "Use this slide while running a live demo or showing screenshots from ingestion, MongoDB, retrieval, and the UI.")
    add_slide_transition(slide)


def create_advantages(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 13)
    add_title(slide, "Advantages", "Why this design is useful")

    cards = [
        ("Reduced Hallucination", "Answers are tied to retrieved PubMed context.", "SAFE"),
        ("Source Visibility", "Users can inspect the chunks used for the answer.", "SRC"),
        ("Context-Aware", "Recent and older memory improve follow-up handling.", "MEM"),
        ("Beginner-Friendly Stack", "Python, MongoDB, sentence-transformers, Gradio.", "PY"),
    ]
    positions = [(0.85, 1.65), (6.85, 1.65), (0.85, 4.00), (6.85, 4.00)]
    for (title, detail, icon), (x, y) in zip(cards, positions):
        add_icon_card(slide, title, detail, x, y, 5.45, 1.55, icon, GREEN if icon in {"SAFE", "MEM"} else CYAN)

    add_notes(slide, "Summarize the main technical benefits, especially grounding and source visibility.")
    add_slide_transition(slide)


def create_future(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 14)
    add_title(slide, "Future Improvements", "Next engineering steps")

    items = [
        ("Gemma or Llama LLM", "Add a local or API-based language model for stronger natural answers."),
        ("Better Vector Search", "Use MongoDB Atlas Vector Search, FAISS, or Qdrant for scalable retrieval."),
        ("Voice Assistant", "Speech-to-text input and text-to-speech output for hands-free use."),
        ("Medical Recommendation Layer", "Only after stronger validation, evaluation, and clinical safety checks."),
    ]
    for idx, (title, detail) in enumerate(items):
        y = 1.55 + idx * 1.18
        add_panel(slide, 1.05, y, 10.9, 0.82)
        add_text(slide, f"{idx + 1}", 1.28, y + 0.22, 0.45, 0.20, 15, CYAN, True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.95, y + 0.14, 3.3, 0.24, 14, WHITE, True)
        add_text(slide, detail, 5.20, y + 0.16, 6.35, 0.38, 11, MUTED)

    add_notes(slide, "Make clear that clinical recommendation systems require validation and safety work beyond this prototype.")
    add_slide_transition(slide)


def create_conclusion(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 15)
    add_title(slide, "Conclusion", "What the project demonstrates")

    add_panel(slide, 1.05, 1.55, 11.15, 4.80, "Summary")
    add_rich_text(
        slide,
        [
            "Built an end-to-end RAG chatbot over PubMed medical abstracts",
            "Ingested CSV data into chunks with embeddings and metadata",
            "Used MongoDB for chunk storage, duplicate prevention, and memory",
            "Retrieved top medical context with cosine similarity",
            "Displayed grounded answers and sources through a Gradio interface",
        ],
        1.55,
        2.18,
        10.15,
        2.85,
        17,
        WHITE,
    )

    add_text(slide, "RAG + MongoDB + Embeddings + Memory = explainable medical QA prototype", 1.80, 5.65, 9.6, 0.32, 16, CYAN, True, align=PP_ALIGN.CENTER)

    add_notes(slide, "Close the technical explanation by tying together ingestion, retrieval, memory, and UI.")
    add_slide_transition(slide)


def create_github(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 16)
    add_title(slide, "GitHub Repository", "Project code, README, and setup instructions")

    add_panel(slide, 1.10, 1.70, 10.95, 3.35, "Repository")
    add_text(slide, REPO_URL, 1.55, 2.55, 9.95, 0.38, 22, CYAN, True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(slide, "Includes ingestion, retrieval, memory, RAG prompt building, Gradio UI, and test questions.", 2.00, 3.32, 9.0, 0.32, 14, WHITE, align=PP_ALIGN.CENTER)
    add_chip(slide, "src/", 3.25, 4.18, 0.72)
    add_chip(slide, "ui/", 4.10, 4.18, 0.70, GREEN)
    add_chip(slide, "app.py", 4.95, 4.18, 0.82)
    add_chip(slide, "README.md", 5.95, 4.18, 1.15, YELLOW)
    add_chip(slide, "test_questions.txt", 7.28, 4.18, 1.65)

    add_panel(slide, 3.20, 5.55, 6.90, 0.55)
    add_text(slide, "Run locally: python app.py", 3.55, 5.73, 6.25, 0.18, 12, WHITE, True, font=FONT_MONO, align=PP_ALIGN.CENTER)

    add_notes(slide, "Show the GitHub link and mention that the project is structured as separate logic and UI modules.")
    add_slide_transition(slide)


def create_thanks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 17)

    add_text(slide, "Thank You", 3.05, 2.35, 7.25, 0.75, 48, WHITE, True, FONT_TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, "Questions?", 4.72, 3.30, 3.90, 0.42, 24, CYAN, True, align=PP_ALIGN.CENTER)
    add_text(slide, PROJECT_TITLE, 4.10, 4.15, 5.05, 0.28, 15, MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, REPO_URL, 3.12, 4.58, 7.0, 0.25, 11, CYAN, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # Minimal AI-medical visual mark.
    center_x, center_y = 6.65, 5.70
    for angle_idx, offset in enumerate([(0, -0.55), (0.50, -0.25), (0.50, 0.35), (0, 0.62), (-0.50, 0.35), (-0.50, -0.25)]):
        x = center_x + offset[0]
        y = center_y + offset[1]
        dot = slide.shapes.add_shape(SHAPE.OVAL, Inches(x), Inches(y), Inches(0.13), Inches(0.13))
        set_fill(dot, CYAN if angle_idx % 2 == 0 else GREEN)
        dot.line.fill.background()

    add_notes(slide, "End with questions and invite the audience to view the repository.")
    add_slide_transition(slide)


def build_deck() -> Presentation:
    """Build the complete 17-slide deck."""

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_cover(prs)
    create_intro(prs)
    create_rag_slide(prs)
    create_project_idea(prs)
    create_dataset(prs)
    create_architecture(prs)
    create_chunk_embeddings(prs)
    create_mongodb_storage(prs)
    create_retrieval(prs)
    create_memory(prs)
    create_gradio_ui(prs)
    create_demo(prs)
    create_advantages(prs)
    create_future(prs)
    create_conclusion(prs)
    create_github(prs)
    create_thanks(prs)

    return prs


def main() -> None:
    """Create the PowerPoint presentation."""

    OUTPUT_DIR.mkdir(exist_ok=True)
    create_visual_assets()
    presentation = build_deck()
    presentation.save(OUTPUT_FILE)
    print(f"Created: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
